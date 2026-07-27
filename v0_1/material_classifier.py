"""
AION Material Classifier.
Classifies uploaded files into: textbook | notes | question_bank | slides | unknown
Uses rule-based filename and content structural signals.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from typing import Optional, List

import fitz

@dataclass
class ClassificationResult:
    material_type: str          # textbook | notes | question_bank | slides | unknown
    confidence: float           # 0.0 to 1.0
    signals: List[str]          # human-readable reasons
    suggested_subject: Optional[str] = None


TEXTBOOK_SIGNALS = [
    r"\b(table\s+of\s+contents|contents)\b",
    r"\b(chapter\s+\d+|unit\s+\d+|module\s+\d+)\b",
    r"\bISBN\b",
    r"\b(edition|published\s+by|mcgraw|pearson|wiley|springer|phi\s+learning|technical\s+publications)\b",
    r"\b(preface|foreword|acknowledgements?)\b",
    r"\b(theorem|lemma|corollary|proof)\b",
    r"\b(bibliography|references|further\s+reading)\b",
]

NOTES_SIGNALS = [
    r"\b(lecture\s+notes?|class\s+notes?|hand\s*written|handout)\b",
    r"\b(module\s+\d+\s+notes?|unit\s+\d+\s+notes?)\b",
    r"\b(prepared\s+by|notes?\s+by|compiled\s+by)\b",
    r"\b(department\s+of|school\s+of|college\s+of)\b",
    r"\b(sem\s*[-:]?\s*\d|semester\s+\d)\b",
    r"^\s*[-•*]\s+.{10,}",
]

QUESTION_BANK_SIGNALS = [
    r"\b(question\s+bank|question\s+paper|previous\s+(year|sem|semester))\b",
    r"\b(model\s+(question|paper)|sample\s+(question|paper))\b",
    r"\b(max(imum)?\s+marks?|total\s+marks?|time\s*:\s*\d+\s*hours?)\b",
    r"\b(part\s*[–-]?\s*[abcABC]|section\s+[abcABC])\b",
    r"^\s*Q[\.\)]\s*\d+",
    r"^\s*\d+[\.\)]\s+.{10,}",
    r"\b(2\s+marks?|5\s+marks?|10\s+marks?|CO\d+|RBT\s*:\s*L\d)\b",
    r"\bVTU\b|\bKSSEM\b|\bRVCE\b|\bBMSCE\b",
]

SLIDES_SIGNALS = [
    r"\b(slide\s+\d+|presentation|presented\s+by)\b",
    r"^\s*\[Slide\s+\d+\]",
    r"\b(click\s+to\s+edit|click\s+to\s+add)\b",
    r"^\s*[•▪▸►]\s+.{5,}",
]

SUBJECT_HINTS = {
    "machine learning": "Machine Learning",
    "deep learning": "Deep Learning",
    "neural network": "Deep Learning",
    "artificial intelligence": "Artificial Intelligence",
    "natural language processing": "NLP",
    "nlp": "NLP",
    "computer vision": "Computer Vision",
    "reinforcement learning": "Reinforcement Learning",
    "data mining": "Data Mining",
    "data structures": "Data Structures & Algorithms",
    "algorithms": "Data Structures & Algorithms",
    "statistics": "Statistics for ML",
    "probability": "Statistics for ML",
    "thermodynamics": "Thermodynamics",
    "fluid mechanics": "Fluid Mechanics",
    "control systems": "Control Systems",
    "digital signal processing": "DSP",
    "vlsi": "VLSI Design",
    "power systems": "Power Systems",
    "structural analysis": "Structural Analysis",
}


def _sample_text(file_path: str, max_pages: int = 8, max_chars: int = 4000) -> str:
    ext = Path(file_path).suffix.lower()
    text = ""

    try:
        if ext == ".pdf":
            doc = fitz.open(file_path)
            parts = []
            pages = list(doc)[:max_pages]
            for page in pages:
                parts.append(page.get_text("text") or "")
            doc.close()
            text = "\n".join(parts)

        elif ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                parts = []
                for slide in list(prs.slides)[:20]:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            parts.append(shape.text)
                text = "\n".join(parts)
            except Exception:
                text = ""

        elif ext in (".txt", ".md"):
            text = Path(file_path).read_text(encoding="utf-8", errors="ignore")

    except Exception:
        text = ""

    return text[:max_chars]


def _count_signal_hits(text: str, patterns: list[str]) -> int:
    hits = 0
    for p in patterns:
        if re.search(p, text, flags=re.I | re.M):
            hits += 1
    return hits


def _infer_subject(text: str, filename: str) -> Optional[str]:
    combined = (filename + " " + text[:1000]).lower()
    for keyword, subject in SUBJECT_HINTS.items():
        if keyword in combined:
            return subject
    return None


def classify_material(file_path: str) -> ClassificationResult:
    """
    Auto-classify an uploaded file into textbook | notes | question_bank | slides | unknown.
    """
    filename = Path(file_path).name.lower()
    ext = Path(file_path).suffix.lower()
    text = _sample_text(file_path)
    combined = filename + "\n" + text

    signals_found = []
    scores = {
        "textbook": 0.0,
        "notes": 0.0,
        "question_bank": 0.0,
        "slides": 0.0,
    }

    if re.search(r"(text\s*book|tb_|_tb|textbook)", filename):
        scores["textbook"] += 0.5
        signals_found.append("filename:textbook")

    if re.search(r"(note|lec|lecture|handout|module_\d|unit_\d)", filename):
        scores["notes"] += 0.5
        signals_found.append("filename:notes")

    if re.search(r"(qbank|question.?bank|qp_|_qp|previous|model.?paper|paper_\d{4})", filename):
        scores["question_bank"] += 0.6
        signals_found.append("filename:question_bank")

    if re.search(r"(slide|ppt|presentation)", filename) or ext in (".pptx", ".ppt"):
        scores["slides"] += 0.5
        signals_found.append("filename/ext:slides")

    tb = _count_signal_hits(combined, TEXTBOOK_SIGNALS)
    nt = _count_signal_hits(combined, NOTES_SIGNALS)
    qb = _count_signal_hits(combined, QUESTION_BANK_SIGNALS)
    sl = _count_signal_hits(combined, SLIDES_SIGNALS)

    scores["textbook"] += tb * 0.1
    scores["notes"] += nt * 0.1
    scores["question_bank"] += qb * 0.15
    scores["slides"] += sl * 0.12

    if tb: signals_found.append(f"content:textbook_signals={tb}")
    if nt: signals_found.append(f"content:notes_signals={nt}")
    if qb: signals_found.append(f"content:qbank_signals={qb}")
    if sl: signals_found.append(f"content:slides_signals={sl}")

    if re.search(r"(table\s+of\s+contents|^\s*contents\s*$)", combined, re.I | re.M):
        scores["textbook"] += 0.3
        signals_found.append("content:TOC_found")

    if ext in (".pptx", ".ppt"):
        scores["slides"] += 0.4

    total = sum(scores.values()) or 1.0
    normalised = {k: v / total for k, v in scores.items()}
    best_type = max(normalised, key=normalised.get)
    best_conf = normalised[best_type]

    if best_conf < 0.30:
        best_type = "unknown"

    subject = _infer_subject(text, filename)

    return ClassificationResult(
        material_type=best_type,
        confidence=round(best_conf, 3),
        signals=signals_found,
        suggested_subject=subject,
    )


def classify_batch(file_paths: list[str]) -> list[ClassificationResult]:
    return [classify_material(fp) for fp in file_paths]
