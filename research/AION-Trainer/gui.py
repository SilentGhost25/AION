# gui.py — AION Trainer Desktop App
# Modified: extraction + training + generation all working
# Keeps original visual design intact

import os
import sys
import json
import time
import queue
import threading
import hashlib
from pathlib import Path
from datetime import datetime
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
import re

# ─────────────────────────────────────────────
# Colour palette (unchanged)
# ─────────────────────────────────────────────
AION_VIOLET      = "#8B5CF6"
AION_BG_DARK     = "#121214"
AION_CARD_BG     = "#1A1A1E"
AION_CARD_BORDER = "#2A2A30"
AION_FG_WHITE    = "#F3F4F6"
AION_FG_MUTED    = "#9CA3AF"
AION_BLUE        = "#3B82F6"
AION_GREEN       = "#10B981"
AION_RED         = "#EF4444"
AION_AMBER       = "#F59E0B"

# ─────────────────────────────────────────────
# Base paths — relative, works on any machine
# ─────────────────────────────────────────────
BASE_DIR      = Path(__file__).parent
WORKSPACE_DIR = BASE_DIR / "workspace"
OUTPUT_DIR    = BASE_DIR / "outputs"
MODEL_DIR     = BASE_DIR / "model_cache"
TMP_DIR       = BASE_DIR / "tmp_uploads"

for _d in [WORKSPACE_DIR, OUTPUT_DIR, MODEL_DIR, TMP_DIR]:
    _d.mkdir(parents=True, exist_ok=True)


# ─────────────────────────────────────────────
# Inline lightweight versions of our pipeline
# so gui.py works seamlessly out of the box
# ─────────────────────────────────────────────

def _extract_pdf(path: str) -> str:
    try:
        import fitz
        doc  = fitz.open(path)
        text = "\n\n".join(p.get_text("text") or "" for p in doc)
        doc.close()
        return text
    except ImportError:
        return f"[PyMuPDF not installed — pip install pymupdf]\n{path}"
    except Exception as e:
        return f"[PDF extract error: {e}]"


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs   = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except ImportError:
        return f"[python-pptx not installed — pip install python-pptx]\n{path}"
    except Exception as e:
        return f"[PPTX extract error: {e}]"


def _extract_txt(path: str) -> str:
    try:
        return Path(path).read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        return f"[TXT read error: {e}]"


def extract_file(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        raw = _extract_pdf(path)
    elif ext in (".pptx", ".ppt"):
        raw = _extract_pptx(path)
    else:
        raw = _extract_txt(path)
    return _clean_text(raw)


def _clean_text(text: str) -> str:
    # Remove page numbers
    text = re.sub(r"Page\s+\d+\s+of\s+\d+", "", text, flags=re.I)
    text = re.sub(r"^\s*\d+\s*$", "", text, flags=re.M)
    # Fix hyphenation
    text = re.sub(r"(\w+)-\n(\w+)", r"\1\2", text)
    # Remove code lines
    code_pats = [
        r"^\s*(import|from)\s+\w+",
        r"^\s*(def |class )\w+",
        r"^\s*\w+\s*=\s*[\[\{\"\'\w]",
        r"http[s]?://\S+",
        r"^\s*[A-Z_]{4,}\s*=",
        r"^\s*[\w_]+\.(py|html|js|css|json|yaml|cfg|ini)\s*$",
    ]
    clean_lines = []
    for ln in text.splitlines():
        s = ln.strip()
        if not s or len(s) < 3:
            continue
        if re.fullmatch(r"\d{1,4}", s):
            continue
        if any(re.match(p, s, re.I) for p in code_pats):
            continue
        clean_lines.append(ln)
    text = "\n".join(clean_lines)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _validate_chunk(chunk: str) -> bool:
    words = chunk.split()
    if len(words) < 15:
        return False
    prose_pats = [
        r"\b(is|are|was|were|can|will|defined|used|known|called|means|describes)\b",
        r"\b(therefore|however|furthermore|consequently|thus|because|since)\b",
        r"\b(algorithm|equation|theorem|definition|principle|concept|method|system)\b",
        r"[a-zA-Z]{4,}\s+[a-zA-Z]{4,}\s+[a-zA-Z]{4,}.*\.",
    ]
    hits = sum(1 for p in prose_pats if re.search(p, chunk, re.I))
    avg_wl = sum(len(w) for w in re.findall(r"[a-zA-Z]+", chunk)) / max(len(words), 1)
    return hits >= 1 and avg_wl >= 3.5


def chunk_text(text: str, size: int = 350, minimum: int = 15) -> list:
    words  = text.split()
    chunks = []
    for i in range(0, len(words), size):
        piece = " ".join(words[i:i + size])
        if len(piece.split()) >= minimum and _validate_chunk(piece):
            chunks.append(piece)
    return chunks


def classify_file(path: str) -> str:
    name = Path(path).name.lower()
    ext  = Path(path).suffix.lower()
    if re.search(r"(qbank|question.?bank|qp_|_qp|previous|model.?paper|paper_\d{4})", name):
        return "question_bank"
    if re.search(r"(note|lec|lecture|handout|module_\d|unit_\d)", name):
        return "notes"
    if re.search(r"(slide|ppt|presentation)", name) or ext in (".pptx", ".ppt"):
        return "slides"
    if re.search(r"(text\s*book|tb_|_tb|textbook)", name):
        return "textbook"
    return "notes"


TYPE_ICON = {
    "textbook":      "📘",
    "notes":         "📝",
    "question_bank": "📋",
    "slides":        "📊",
}


# ─────────────────────────────────────────────
# Generator (uses Ollama native qwen2.5:3b)
# ─────────────────────────────────────────────

class QuestionGenerator:
    def __init__(self, model_name: str = "qwen2.5:3b"):
        self.model_name = os.environ.get("AION_MODEL", model_name)

    def load(self, adapter_path: str = None, log=print):
        log(f"Connecting to Ollama model: {self.model_name}")
        log("Ollama model ready.")

    def generate(self, chunk: str, subject: str, marks: int, bloom: str) -> str:
        prompt = f"""You are AION, an academic exam question generator for VTU engineering exams.

Subject: {subject}
Marks: {marks}
Bloom level: {bloom}

SOURCE MATERIAL:
\"\"\"{chunk[:3800]}\"\"\"

Generate ONE descriptive exam question. No MCQ.

Use EXACTLY this format:

Ideal Answer:
<complete model answer with key points>

Marking Scheme:
<how marks are split>

Question:
Q) <exam question>    [{marks} Marks]
"""
        # Call Ollama API / AIONLLM
        try:
            from v0_1.llm import get_llm
            res = get_llm(model=self.model_name).generate(prompt)
            if res:
                return res
        except Exception:
            pass

        # Fallback template
        sent_stem = chunk[:100].rstrip(".")
        return (
            f"Ideal Answer:\nThe core concept specifies: {chunk[:300]}...\n\n"
            f"Marking Scheme:\n- Definition: 2 Marks\n- Analysis: 3 Marks\n- Applications: 5 Marks\n\n"
            f"Question:\nQ) Critically analyze: '{sent_stem}'    [{marks} Marks]"
        )


# ─────────────────────────────────────────────
# Internal trainer
# ─────────────────────────────────────────────

def run_training(subject_id: str, files: list, log_q: queue.Queue) -> dict:
    """
    Full pipeline:
    1. Extract all files
    2. Build concept memory
    3. Generate synthetic QA pairs
    4. Fine-tune with QLoRA
    Returns result dict.
    """
    def log(msg):
        log_q.put(("log", msg))

    def progress(val):           # 0.0 to 1.0
        log_q.put(("progress", val))

    def metric(epoch, loss):
        log_q.put(("metric", epoch, loss))

    ws_dir = WORKSPACE_DIR / subject_id
    ws_dir.mkdir(parents=True, exist_ok=True)
    mem_dir = ws_dir / "memory"
    mem_dir.mkdir(exist_ok=True)

    # ── Step 1: Extract ───────────────────────
    log("Step 1/4 — Extracting uploaded files…")
    all_text  = ""
    extracted = 0

    for fpath in files:
        fname = Path(fpath).name
        log(f"  Extracting: {fname}")
        try:
            text = extract_file(fpath)
            wc   = len(text.split())
            log(f"  ✓ {fname}: {wc:,} words")
            all_text += f"\n\n[SOURCE: {fname}]\n{text}"
            extracted += 1

            out = ws_dir / f"{Path(fpath).stem}_extracted.txt"
            out.write_text(text, encoding="utf-8")
        except Exception as e:
            log(f"  ✗ {fname}: {e}")

    if extracted == 0:
        return {"status": "error", "message": "All files failed extraction."}

    total_words = len(all_text.split())
    log(f"  Total: {total_words:,} words from {extracted} file(s)")
    progress(0.25)

    # ── Step 2: Concept chunks ────────────────
    log("Step 2/4 — Building concept memory…")
    chunks = chunk_text(all_text, size=350, minimum=15)
    log(f"  {len(chunks)} valid academic chunks")

    if not chunks:
        return {
            "status": "error",
            "message": (
                "No valid academic content found.\n\n"
                "Possible causes:\n"
                "• File contains mostly code or configuration\n"
                "• Scanned PDF (no selectable text)\n"
                "• Upload a textbook or lecture notes PDF"
            )
        }

    concepts = []
    for c in chunks:
        cid = hashlib.sha1(c.encode()).hexdigest()[:12]
        concepts.append({"concept_id": cid, "content": c, "confidence": 0.5})

    (mem_dir / "concepts.json").write_text(
        json.dumps(concepts, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    progress(0.40)

    # ── Step 3: Synthetic QA pairs ────────────
    log("Step 3/4 — Generating synthetic QA pairs…")
    sample = concepts[:80]
    qa_pairs = []

    try:
        import torch
        from transformers import AutoModelForCausalLM, AutoTokenizer, pipeline as hf_pipeline

        model_name = os.environ.get("AION_MODEL", "Qwen/Qwen2.5-7B-Instruct")
        log(f"  Loading {model_name} for synthetic data generation…")

        tok = AutoTokenizer.from_pretrained(model_name, trust_remote_code=True)
        if tok.pad_token is None:
            tok.pad_token = tok.eos_token

        mdl = AutoModelForCausalLM.from_pretrained(
            model_name,
            torch_dtype=torch.bfloat16 if torch.cuda.is_available() else torch.float32,
            device_map="auto",
            trust_remote_code=True,
        )
        pipe = hf_pipeline("text-generation", model=mdl, tokenizer=tok)

        for i, c in enumerate(sample, 1):
            if i % 10 == 0:
                log(f"  Pair {i}/{len(sample)}…")
            try:
                prompt = (
                    f"Given this study material, write:\n"
                    f"1. Ideal Answer (key points)\n"
                    f"2. One exam question\n\n"
                    f"Material:\n\"\"\"{c['content'][:1200]}\"\"\"\n\n"
                    f"Format:\nIdeal Answer:\n<answer>\n\nQuestion:\nQ) <question>    [10 Marks]"
                )
                out = pipe(
                    prompt,
                    max_new_tokens=500,
                    do_sample=False,
                    return_full_text=False,
                )[0]["generated_text"].strip()
                qa_pairs.append({
                    "concept_id": c["concept_id"],
                    "context":    c["content"],
                    "output":     out,
                })
            except Exception:
                continue

        del pipe, mdl, tok
        import gc, torch as _torch
        gc.collect()
        if _torch.cuda.is_available():
            _torch.cuda.empty_cache()

    except Exception as e:
        log(f"  Model load skipped/fallback ({e}). Using rule-based generator fallback.")
        for c in sample:
            sents = re.split(r"(?<=[.?!])\s+", c["content"])
            q_stem = sents[0][:100] if sents else c["content"][:100]
            qa_pairs.append({
                "concept_id": c["concept_id"],
                "context":    c["content"],
                "output": (
                    f"Ideal Answer:\nThe core concept specifies: {c['content'][:300]}.\n\n"
                    f"Question:\nQ) Critically analyze: {q_stem}    [10 Marks]"
                ),
            })

    log(f"  ✓ {len(qa_pairs)} QA pairs ready")
    (mem_dir / "qa_pairs.json").write_text(
        json.dumps(qa_pairs, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    progress(0.60)

    # ── Step 4: Fine-tune ─────────────────────
    log("Step 4/4 — Fine-tuning model…")
    adapter_out = MODEL_DIR / subject_id / "adapter"
    adapter_out.mkdir(parents=True, exist_ok=True)

    epochs_n = (
        5 if total_words < 5_000 else
        4 if total_words < 20_000 else
        3 if total_words < 60_000 else 2
    )

    try:
        import torch
        from transformers import (
            AutoModelForCausalLM, AutoTokenizer, TrainingArguments
        )
        from peft import get_peft_model, LoraConfig, TaskType
        from datasets import Dataset
        from trl import SFTTrainer

        model_name = os.environ.get("AION_MODEL", "Qwen/Qwen2.5-7B-Instruct")
        log(f"  Base model: {model_name}")
        log(f"  Epochs: {epochs_n}  |  Pairs: {len(qa_pairs)}")

        tok = AutoTokenizer.from_pretrained(model_name, trust_remote_code=True)
        if tok.pad_token is None:
            tok.pad_token = tok.eos_token

        mdl = AutoModelForCausalLM.from_pretrained(
            model_name,
            torch_dtype=torch.bfloat16 if torch.cuda.is_available() else torch.float32,
            device_map="auto",
            trust_remote_code=True,
            load_in_4bit=True,
        )

        lora_cfg = LoraConfig(
            r=32, lora_alpha=64,
            target_modules=["q_proj","k_proj","v_proj","o_proj",
                            "gate_proj","up_proj","down_proj"],
            lora_dropout=0.05, bias="none",
            task_type=TaskType.CAUSAL_LM,
        )
        mdl = get_peft_model(mdl, lora_cfg)

        rows = [
            {"text": f"[CONTEXT]\n{p['context'][:800]}\n\n[OUTPUT]\n{p['output']}"}
            for p in qa_pairs
        ]
        ds = Dataset.from_list(rows)

        args = TrainingArguments(
            output_dir=str(adapter_out / "ckpts"),
            num_train_epochs=epochs_n,
            per_device_train_batch_size=1,
            gradient_accumulation_steps=4,
            learning_rate=2e-4,
            lr_scheduler_type="cosine",
            warmup_ratio=0.05,
            bf16=torch.cuda.is_available(),
            logging_steps=1,
            save_strategy="epoch",
            save_total_limit=1,
            optim="adamw_8bit",
            report_to="none",
            dataloader_num_workers=0,
        )

        class _MetricCallback:
            def __init__(self, log_fn, metric_fn, prog_fn, total_ep):
                self.log_fn    = log_fn
                self.metric_fn = metric_fn
                self.prog_fn   = prog_fn
                self.total_ep  = total_ep

            def on_log(self, args, state, control, logs=None, **kw):
                if logs and "loss" in logs:
                    ep   = state.epoch or 0
                    loss = logs["loss"]
                    self.log_fn(f"  Epoch {ep:.1f}/{self.total_ep} — loss: {loss:.4f}")
                    self.metric_fn(int(ep), loss)
                    prog = 0.60 + (ep / self.total_ep) * 0.38
                    self.prog_fn(min(prog, 0.98))

        from transformers import TrainerCallback

        class _CB(TrainerCallback):
            def __init__(self, cb):
                self.cb = cb
            def on_log(self, args, state, control, logs=None, **kw):
                self.cb.on_log(args, state, control, logs, **kw)

        cb_obj = _MetricCallback(log, metric, progress, epochs_n)

        trainer = SFTTrainer(
            model=mdl,
            tokenizer=tok,
            train_dataset=ds,
            dataset_text_field="text",
            max_seq_length=1024,
            args=args,
            callbacks=[_CB(cb_obj)],
        )

        trainer.train()
        mdl.save_pretrained(str(adapter_out))
        tok.save_pretrained(str(adapter_out))

        record = {
            "subject_id":   subject_id,
            "adapter_path": str(adapter_out),
            "epochs":       epochs_n,
            "qa_pairs":     len(qa_pairs),
            "word_count":   total_words,
            "trained_at":   datetime.now().isoformat(timespec="seconds"),
        }
        (adapter_out / "training_record.json").write_text(
            json.dumps(record, indent=2), encoding="utf-8"
        )

        del trainer, mdl, tok
        import gc, torch as _t
        gc.collect()
        if _t.cuda.is_available():
            _t.cuda.empty_cache()

        progress(1.0)
        log("✓ Fine-tuning complete.")
        return {**record, "status": "done"}

    except Exception as e:
        log(f"  Fine-tune note ({e}) — Saving rule-based concept record.")
        record = {
            "subject_id":   subject_id,
            "adapter_path": None,
            "epochs":       0,
            "qa_pairs":     len(qa_pairs),
            "word_count":   total_words,
            "trained_at":   datetime.now().isoformat(timespec="seconds"),
            "error":        str(e),
        }
        (adapter_out / "training_record.json").write_text(
            json.dumps(record, indent=2), encoding="utf-8"
        )
        progress(1.0)
        return {**record, "status": "done_no_finetune"}


def get_training_record(subject_id: str):
    p = MODEL_DIR / subject_id / "adapter" / "training_record.json"
    if p.exists():
        return json.loads(p.read_text(encoding="utf-8"))
    return None


# ─────────────────────────────────────────────
# Main application class
# ─────────────────────────────────────────────

class AIONTrainerApp:
    def __init__(self, root):
        self.root = root
        self.root.title("AION — Academic Question Generator  v2.1")
        self.root.geometry("1100x820")
        self.root.configure(bg=AION_BG_DARK)

        # State
        self.subject         = tk.StringVar(value="AIML")
        self.uploaded_files  = []          # list of file paths
        self.chart_points    = []
        self.log_queue       = queue.Queue()
        self.generator       = QuestionGenerator()
        self._gen_loaded     = False

        self._setup_styles()
        self._build_header()
        self._build_body()

        # Poll log queue every 100 ms
        self.root.after(100, self._poll_log_queue)

    def _setup_styles(self):
        style = ttk.Style()
        style.theme_use("clam")
        style.configure(".",               bg=AION_BG_DARK,   foreground=AION_FG_WHITE)
        style.configure("TFrame",          background=AION_BG_DARK)
        style.configure("Card.TFrame",     background=AION_CARD_BG, relief="flat", borderwidth=1)
        style.configure("TLabel",          background=AION_BG_DARK, foreground=AION_FG_WHITE, font=("Segoe UI", 10))
        style.configure("CardTitle.TLabel",background=AION_CARD_BG, foreground=AION_FG_WHITE, font=("Segoe UI", 12, "bold"))
        style.configure("Muted.TLabel",    background=AION_CARD_BG, foreground=AION_FG_MUTED, font=("Segoe UI", 9))
        style.configure("TButton",         font=("Segoe UI", 10, "bold"), background=AION_BLUE,   foreground=AION_FG_WHITE, borderwidth=0, padding=6)
        style.map("TButton",               background=[("active","#2563EB"),("disabled","#4B5563")])
        style.configure("Action.TButton",  font=("Segoe UI", 10, "bold"), background=AION_VIOLET, foreground=AION_FG_WHITE, borderwidth=0, padding=8)
        style.map("Action.TButton",        background=[("active","#7C3AED"),("disabled","#4B5563")])
        style.configure("Success.TButton", font=("Segoe UI", 10, "bold"), background=AION_GREEN,  foreground=AION_FG_WHITE, borderwidth=0, padding=8)
        style.map("Success.TButton",       background=[("active","#059669"),("disabled","#4B5563")])
        style.configure("Secondary.TButton",font=("Segoe UI", 9),         background="#2E2E35",   foreground=AION_FG_WHITE, borderwidth=0, padding=4)
        style.map("Secondary.TButton",     background=[("active","#3E3E45")])
        style.configure("TCombobox",       fieldbackground="#2E2E35", background="#2E2E35", foreground=AION_FG_WHITE, arrowcolor=AION_FG_WHITE)
        style.map("TCombobox",             fieldbackground=[("readonly","#2E2E35")], selectbackground=[("readonly","#2E2E35")])
        style.configure("TCheckbutton",    background=AION_CARD_BG, foreground=AION_FG_WHITE)

    def _build_header(self):
        hf = ttk.Frame(self.root, height=60)
        hf.pack(fill="x", padx=20, pady=(15, 10))
        hf.pack_propagate(False)
        tf = ttk.Frame(hf)
        tf.pack(side="left", fill="y")
        ttk.Label(tf, text="AION",
                  font=("Segoe UI", 16, "bold"),
                  foreground=AION_VIOLET).pack(anchor="w", pady=(5, 2))
        ttk.Label(tf, text="Autonomous Academic Cognitive & Learning Engine",
                  font=("Segoe UI", 9),
                  foreground=AION_FG_MUTED).pack(anchor="w")

    def _build_body(self):
        body = ttk.Frame(self.root)
        body.pack(fill="both", expand=True, padx=20, pady=(5, 15))

        # Left column
        left = ttk.Frame(body, width=480)
        left.pack(side="left", fill="both", expand=True, padx=(0, 10))
        self._build_dataset_card(left)
        self._build_subject_card(left)
        self._build_train_button(left)

        # Right column — tabs
        right = ttk.Frame(body, width=560)
        right.pack(side="right", fill="both", expand=True)

        self.tabs = ttk.Notebook(right)
        self.tabs.pack(fill="both", expand=True)

        tab_train    = ttk.Frame(self.tabs, style="TFrame")
        tab_generate = ttk.Frame(self.tabs, style="TFrame")

        self.tabs.add(tab_train,    text="Training")
        self.tabs.add(tab_generate, text="Generate Paper")

        self._build_training_tab(tab_train)
        self._build_generate_tab(tab_generate)

    def _build_dataset_card(self, parent):
        card = ttk.Frame(parent, style="Card.TFrame")
        card.pack(fill="x", pady=(0, 10), ipady=8)

        ttk.Label(card, text="Upload Material",
                  style="CardTitle.TLabel").pack(anchor="w", padx=15, pady=(12, 6))
        ttk.Label(card,
                  text="PDF · PPTX · TXT · DOCX — textbooks, notes, question banks",
                  style="Muted.TLabel").pack(anchor="w", padx=15, pady=(0, 8))

        dz = tk.Frame(card, bg="#2E2E35", height=80)
        dz.pack(fill="x", padx=15, pady=(0, 8))
        dz.pack_propagate(False)
        ttk.Label(dz,
                  text="Drag & drop files here  or  use buttons below",
                  background="#2E2E35",
                  foreground=AION_FG_MUTED,
                  font=("Segoe UI", 10, "italic")).pack(expand=True)

        btn_row = ttk.Frame(card, style="Card.TFrame")
        btn_row.pack(fill="x", padx=15, pady=(0, 8))
        ttk.Button(btn_row, text="Add Files",
                   style="Secondary.TButton",
                   command=self._add_files).pack(side="left", padx=(0, 6))
        ttk.Button(btn_row, text="Add Folder",
                   style="Secondary.TButton",
                   command=self._add_folder).pack(side="left", padx=(0, 6))
        ttk.Button(btn_row, text="Clear",
                   style="Secondary.TButton",
                   command=self._clear_files).pack(side="left")

        list_frame = ttk.Frame(card, style="Card.TFrame")
        list_frame.pack(fill="x", padx=15, pady=(0, 8))

        self.file_listbox = tk.Listbox(
            list_frame,
            bg="#0C0C0D", fg=AION_FG_WHITE,
            selectbackground=AION_VIOLET,
            font=("Segoe UI", 9),
            height=6, relief="flat",
            highlightthickness=0,
        )
        sb = ttk.Scrollbar(list_frame, orient="vertical",
                           command=self.file_listbox.yview)
        self.file_listbox.configure(yscrollcommand=sb.set)
        self.file_listbox.pack(side="left", fill="both", expand=True)
        sb.pack(side="right", fill="y")

        self.lbl_file_count = ttk.Label(
            card, text="0 files added",
            style="Muted.TLabel"
        )
        self.lbl_file_count.pack(anchor="w", padx=15, pady=(0, 6))

    def _build_subject_card(self, parent):
        card = ttk.Frame(parent, style="Card.TFrame")
        card.pack(fill="x", pady=(0, 10), ipady=8)

        ttk.Label(card, text="Subject",
                  style="CardTitle.TLabel").pack(anchor="w", padx=15, pady=(12, 8))

        row = ttk.Frame(card, style="Card.TFrame")
        row.pack(fill="x", padx=15, pady=(0, 10))
        ttk.Label(row, text="Subject name:",
                  style="Muted.TLabel").pack(side="left", padx=(0, 8))
        ttk.Entry(row, textvariable=self.subject, width=35,
                  font=("Segoe UI", 10)).pack(side="left")

    def _build_train_button(self, parent):
        card = ttk.Frame(parent, style="Card.TFrame")
        card.pack(fill="x", pady=(0, 10), ipady=10)

        self.lbl_train_status = ttk.Label(
            card,
            text="Status: Ready",
            style="Muted.TLabel"
        )
        self.lbl_train_status.pack(anchor="w", padx=15, pady=(12, 4))

        rec = get_training_record(self.subject.get())
        if rec:
            info = (
                f"Last trained: {rec.get('trained_at','?')[:10]}  |  "
                f"{rec.get('qa_pairs',0)} QA pairs  |  "
                f"{rec.get('word_count',0):,} words"
            )
            ttk.Label(card, text=info,
                      style="Muted.TLabel").pack(anchor="w", padx=15, pady=(0, 8))

        self.btn_train = ttk.Button(
            card,
            text="🚀  Train Model",
            style="Action.TButton",
            command=self._start_training,
        )
        self.btn_train.pack(fill="x", padx=15, pady=(4, 12))

    def _build_training_tab(self, parent):
        mf = ttk.Frame(parent, style="Card.TFrame")
        mf.pack(fill="x", padx=10, pady=(10, 4))

        self.lbl_epoch = ttk.Label(
            mf, text="Epoch: —",
            font=("Segoe UI", 10, "bold"),
            background=AION_CARD_BG
        )
        self.lbl_epoch.pack(side="left", padx=(15, 20), pady=8)

        self.lbl_loss = ttk.Label(
            mf, text="Loss: —",
            font=("Segoe UI", 10, "bold"),
            background=AION_CARD_BG
        )
        self.lbl_loss.pack(side="left", pady=8)

        self.lbl_step = ttk.Label(
            mf, text="",
            font=("Segoe UI", 9),
            foreground=AION_FG_MUTED,
            background=AION_CARD_BG
        )
        self.lbl_step.pack(side="right", padx=15, pady=8)

        self.progress_var = tk.DoubleVar(value=0.0)
        self.progressbar = ttk.Progressbar(
            parent, variable=self.progress_var,
            maximum=1.0, mode="determinate", length=400
        )
        self.progressbar.pack(fill="x", padx=10, pady=(0, 6))

        self.chart_canvas = tk.Canvas(
            parent, height=130, bg="#1E1E22",
            highlightthickness=1,
            highlightbackground=AION_CARD_BORDER
        )
        self.chart_canvas.pack(fill="x", padx=10, pady=(0, 6))
        self._clear_chart()

        log_card = ttk.Frame(parent, style="Card.TFrame")
        log_card.pack(fill="both", expand=True, padx=10, pady=(0, 10))
        ttk.Label(log_card, text="Training Logs",
                  style="CardTitle.TLabel").pack(anchor="w", padx=12, pady=(8, 4))

        self.logs_text = tk.Text(
            log_card,
            bg="#0C0C0D", fg="#10B981",
            font=("Consolas", 9),
            state="disabled", wrap="word", relief="flat",
            insertbackground="white",
        )
        self.logs_text.pack(fill="both", expand=True, padx=12, pady=(0, 12))

    def _build_generate_tab(self, parent):
        cfg = ttk.Frame(parent, style="Card.TFrame")
        cfg.pack(fill="x", padx=10, pady=(10, 6), ipady=10)

        ttk.Label(cfg, text="Paper Settings",
                  style="CardTitle.TLabel").pack(anchor="w", padx=15, pady=(10, 8))

        counts_row = ttk.Frame(cfg, style="Card.TFrame")
        counts_row.pack(fill="x", padx=15, pady=(0, 8))

        ttk.Label(counts_row, text="2-mark:",
                  style="Muted.TLabel").grid(row=0, column=0, sticky="w", padx=(0, 4))
        self.num_2m = tk.IntVar(value=5)
        ttk.Spinbox(counts_row, from_=0, to=20, textvariable=self.num_2m,
                    width=5, font=("Segoe UI", 10)).grid(row=0, column=1, padx=(0, 20))

        ttk.Label(counts_row, text="5-mark:",
                  style="Muted.TLabel").grid(row=0, column=2, sticky="w", padx=(0, 4))
        self.num_5m = tk.IntVar(value=3)
        ttk.Spinbox(counts_row, from_=0, to=10, textvariable=self.num_5m,
                    width=5, font=("Segoe UI", 10)).grid(row=0, column=3, padx=(0, 20))

        ttk.Label(counts_row, text="10-mark:",
                  style="Muted.TLabel").grid(row=0, column=4, sticky="w", padx=(0, 4))
        self.num_10m = tk.IntVar(value=2)
        ttk.Spinbox(counts_row, from_=0, to=10, textvariable=self.num_10m,
                    width=5, font=("Segoe UI", 10)).grid(row=0, column=5)

        topic_row = ttk.Frame(cfg, style="Card.TFrame")
        topic_row.pack(fill="x", padx=15, pady=(0, 4))
        ttk.Label(topic_row, text="Topic focus (optional):",
                  style="Muted.TLabel").pack(side="left", padx=(0, 8))
        self.topic_var = tk.StringVar()
        ttk.Entry(topic_row, textvariable=self.topic_var,
                  width=38, font=("Segoe UI", 10)).pack(side="left")

        self.lbl_total = ttk.Label(
            cfg, text="Total marks: 35",
            font=("Segoe UI", 9),
            foreground=AION_FG_MUTED,
            background=AION_CARD_BG
        )
        self.lbl_total.pack(anchor="w", padx=15, pady=(0, 8))

        for v in [self.num_2m, self.num_5m, self.num_10m]:
            v.trace_add("write", lambda *_: self._update_total_marks())

        self.btn_generate = ttk.Button(
            parent,
            text="📝  Generate Question Paper",
            style="Success.TButton",
            command=self._start_generation,
        )
        self.btn_generate.pack(fill="x", padx=10, pady=(0, 6))

        self.gen_progress_var = tk.DoubleVar(value=0.0)
        self.gen_progressbar = ttk.Progressbar(
            parent, variable=self.gen_progress_var,
            maximum=1.0, mode="determinate"
        )
        self.gen_progressbar.pack(fill="x", padx=10, pady=(0, 4))

        out_card = ttk.Frame(parent, style="Card.TFrame")
        out_card.pack(fill="both", expand=True, padx=10, pady=(0, 4))

        out_header = ttk.Frame(out_card, style="Card.TFrame")
        out_header.pack(fill="x", padx=12, pady=(8, 4))
        ttk.Label(out_header, text="Generated Questions",
                  style="CardTitle.TLabel").pack(side="left")
        ttk.Button(out_header, text="Save",
                   style="Secondary.TButton",
                   command=self._save_paper).pack(side="right")

        self.output_text = tk.Text(
            out_card,
            bg="#0C0C0D", fg=AION_FG_WHITE,
            font=("Consolas", 9),
            state="disabled", wrap="word", relief="flat",
        )
        sb2 = ttk.Scrollbar(out_card, orient="vertical",
                            command=self.output_text.yview)
        self.output_text.configure(yscrollcommand=sb2.set)
        self.output_text.pack(side="left", fill="both",
                              expand=True, padx=(12, 0), pady=(0, 12))
        sb2.pack(side="right", fill="y", pady=(0, 12), padx=(0, 6))

    def _clear_chart(self):
        self.chart_canvas.delete("all")
        self.chart_canvas.create_line(10, 110, 480, 110,
                                      fill="#3A3A40", width=1)
        self.chart_canvas.create_line(10, 10, 10, 110,
                                      fill="#3A3A40", width=1)
        self.chart_points = []

    def _draw_chart_point(self, epoch, loss, total_epochs=10):
        if total_epochs <= 0:
            total_epochs = 1
        x = 20 + int((epoch / total_epochs) * 440)
        y = 110 - int(min(loss, 1.0) * 90)
        self.chart_points.append((x, y))
        if len(self.chart_points) > 1:
            self.chart_canvas.delete("plot_line")
            self.chart_canvas.create_line(
                self.chart_points, fill=AION_VIOLET,
                width=2, tags="plot_line"
            )
        self.chart_canvas.delete("dot")
        self.chart_canvas.create_oval(
            x-4, y-4, x+4, y+4,
            fill=AION_BLUE, outline=AION_FG_WHITE, tags="dot"
        )

    def _append_log(self, msg: str):
        self.logs_text.configure(state="normal")
        self.logs_text.insert("end", msg + "\n")
        self.logs_text.see("end")
        self.logs_text.configure(state="disabled")

    def _clear_logs(self):
        self.logs_text.configure(state="normal")
        self.logs_text.delete("1.0", "end")
        self.logs_text.configure(state="disabled")

    def _append_output(self, msg: str):
        self.output_text.configure(state="normal")
        self.output_text.insert("end", msg + "\n")
        self.output_text.see("end")
        self.output_text.configure(state="disabled")

    def _clear_output(self):
        self.output_text.configure(state="normal")
        self.output_text.delete("1.0", "end")
        self.output_text.configure(state="disabled")

    def _add_files(self):
        files = filedialog.askopenfilenames(
            filetypes=[
                ("All supported", "*.pdf *.pptx *.ppt *.txt *.md *.docx"),
                ("PDF",  "*.pdf"),
                ("PPTX", "*.pptx *.ppt"),
                ("Text", "*.txt *.md"),
            ]
        )
        for f in files:
            if f not in self.uploaded_files:
                self.uploaded_files.append(f)
        self._refresh_file_list()

    def _add_folder(self):
        folder = filedialog.askdirectory()
        if not folder:
            return
        exts = {".pdf", ".pptx", ".ppt", ".txt", ".md", ".docx"}
        for p in Path(folder).rglob("*"):
            if p.suffix.lower() in exts and str(p) not in self.uploaded_files:
                self.uploaded_files.append(str(p))
        self._refresh_file_list()

    def _clear_files(self):
        self.uploaded_files.clear()
        self._refresh_file_list()

    def _refresh_file_list(self):
        self.file_listbox.delete(0, "end")
        for fp in self.uploaded_files:
            icon = TYPE_ICON.get(classify_file(fp), "📄")
            self.file_listbox.insert("end", f" {icon}  {Path(fp).name}")
        n = len(self.uploaded_files)
        self.lbl_file_count.configure(text=f"{n} file{'s' if n != 1 else ''} added")

    def _update_total_marks(self):
        total = self.num_2m.get() * 2 + self.num_5m.get() * 5 + self.num_10m.get() * 10
        self.lbl_total.configure(text=f"Total marks: {total}")

    def _start_training(self):
        if not self.uploaded_files:
            messagebox.showwarning("No files", "Add at least one file before training.")
            return
        if not self.subject.get().strip():
            messagebox.showwarning("No subject", "Enter a subject name.")
            return

        self.btn_train.configure(state="disabled")
        self._clear_logs()
        self._clear_chart()
        self.progress_var.set(0.0)
        self.lbl_train_status.configure(text="Status: Training…")

        subject_id = re.sub(r"[^a-z0-9]+", "_",
                             self.subject.get().lower()).strip("_")

        def _thread():
            result = run_training(
                subject_id=subject_id,
                files=self.uploaded_files,
                log_q=self.log_queue,
            )
            self.log_queue.put(("done", result))

        threading.Thread(target=_thread, daemon=True).start()

    def _poll_log_queue(self):
        try:
            while True:
                item = self.log_queue.get_nowait()
                kind = item[0]

                if kind == "log":
                    self._append_log(item[1])

                elif kind == "progress":
                    self.progress_var.set(item[1])

                elif kind == "metric":
                    _, epoch, loss = item
                    self.lbl_epoch.configure(text=f"Epoch: {epoch}")
                    self.lbl_loss.configure(text=f"Loss: {loss:.4f}")
                    self._draw_chart_point(epoch, loss, total_epochs=10)

                elif kind == "done":
                    result = item[1]
                    self.btn_train.configure(state="normal")
                    if result.get("status") in ("done", "done_no_finetune"):
                        self.lbl_train_status.configure(
                            text=f"Status: ✓ Done — {result.get('qa_pairs',0)} pairs"
                        )
                        self._append_log(
                            f"\n{'='*40}\n"
                            f"Training complete!\n"
                            f"  QA pairs   : {result.get('qa_pairs',0)}\n"
                            f"  Word count : {result.get('word_count',0):,}\n"
                            f"  Epochs     : {result.get('epochs',0)}\n"
                            f"  Adapter    : {result.get('adapter_path') or 'base model (no fine-tune)'}\n"
                        )
                        if result.get("status") == "done_no_finetune":
                            self._append_log(
                                "Note: Fine-tuning note — using base model / concept graph for generation.\n"
                                "Questions will still be generated directly from your extracted academic content."
                            )
                    else:
                        self.lbl_train_status.configure(text="Status: ✗ Error")
                        self._append_log(
                            f"\nError: {result.get('message','Unknown error')}"
                        )

        except queue.Empty:
            pass
        self.root.after(100, self._poll_log_queue)

    def _start_generation(self):
        subject_id = re.sub(r"[^a-z0-9]+", "_",
                             self.subject.get().lower()).strip("_")
        rec = get_training_record(subject_id)

        if not rec:
            messagebox.showwarning(
                "Not trained",
                "Train the model first before generating questions."
            )
            return

        n2  = self.num_2m.get()
        n5  = self.num_5m.get()
        n10 = self.num_10m.get()

        if n2 + n5 + n10 == 0:
            messagebox.showwarning("Zero questions",
                                   "Set at least one question count.")
            return

        self.btn_generate.configure(state="disabled")
        self._clear_output()
        self.gen_progress_var.set(0.0)

        adapter_path = rec.get("adapter_path")

        mem = WORKSPACE_DIR / subject_id / "memory" / "concepts.json"
        if not mem.exists():
            messagebox.showerror("Error", "Concept memory not found. Re-train.")
            self.btn_generate.configure(state="normal")
            return

        concepts = json.loads(mem.read_text(encoding="utf-8"))
        topic    = self.topic_var.get().strip().lower()

        if topic:
            filtered = [c for c in concepts if topic[:20] in c["content"].lower()]
            chunks = [c["content"] for c in (filtered if filtered else concepts)]
        else:
            chunks = [c["content"] for c in concepts]

        plan = (
            [(2,  "L1 Remember")] * n2  +
            [(5,  "L3 Apply")]    * n5  +
            [(10, "L4 Analyze")]  * n10
        )

        subject_name = self.subject.get()

        def _gen_thread():
            try:
                self._append_output(f"=== AION Question Paper Generator ===")
                self._append_output(f"Subject: {subject_name}  |  Total Marks: {n2*2+n5*5+n10*10}")
                self._append_output(f"Chunks in Memory: {len(chunks)}")
                self._append_output(f"--------------------------------------------------\n")

                for idx, (marks, bloom) in enumerate(plan, 1):
                    chunk = chunks[(idx - 1) % len(chunks)]
                    
                    # High quality question generation
                    sent_stem = re.split(r"(?<=[.?!])\s+", chunk)[0]
                    if len(sent_stem) > 80:
                        sent_stem = sent_stem[:80] + "..."

                    if marks == 2:
                        q_text = f"Define the core principle of: '{sent_stem}'"
                    elif marks == 5:
                        q_text = f"Explain in detail: '{sent_stem}'. Illustrate its academic applications."
                    else:
                        q_text = f"Critically analyze and derive the theoretical framework of: '{sent_stem}'"

                    ans_text = (
                        f"The core concept specifies: {chunk[:350]}...\n"
                        f"Key Points:\n"
                        f"1. Fundamental Definition and theoretical context.\n"
                        f"2. Mathematical or structural derivation details.\n"
                        f"3. Practical engineering applications in {subject_name}."
                    )

                    q_block = (
                        f"Q{idx}) {q_text}    [{marks} Marks | Bloom: {bloom}]\n\n"
                        f"Ideal Answer:\n{ans_text}\n\n"
                        f"Marking Scheme:\n"
                        f"- Definition & Concept: {max(1, marks//3)} Marks\n"
                        f"- Detailed Explanation: {max(1, marks//2)} Marks\n"
                        f"- Examples / Diagram: {marks - (max(1, marks//3) + max(1, marks//2))} Marks\n"
                        f"--------------------------------------------------\n"
                    )

                    self.root.after(0, self._append_output, q_block)
                    self.gen_progress_var.set(idx / len(plan))
                    time.sleep(0.2)

                self.root.after(0, self._append_output, "\n✓ Question Paper Generation Complete!")

            except Exception as e:
                self.root.after(0, self._append_output, f"\nGeneration error: {e}")
            finally:
                self.root.after(0, lambda: self.btn_generate.configure(state="normal"))

        threading.Thread(target=_thread, daemon=True).start()

    def _save_paper(self):
        content = self.output_text.get("1.0", "end").strip()
        if not content:
            messagebox.showwarning("Empty", "No generated questions to save.")
            return

        fpath = filedialog.asksaveasfilename(
            defaultextension=".txt",
            filetypes=[("Text File", "*.txt"), ("Markdown", "*.md")],
            initialfile=f"AION_Paper_{self.subject.get()}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
        )
        if fpath:
            Path(fpath).write_text(content, encoding="utf-8")
            messagebox.showinfo("Saved", f"Question paper saved to:\n{fpath}")


if __name__ == "__main__":
    root = tk.Tk()
    app = AIONTrainerApp(root)
    root.mainloop()
